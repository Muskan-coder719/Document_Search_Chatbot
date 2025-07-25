from drive_auth_test11 import authenticate_drive, fetch_documents
from sentence_transformers import SentenceTransformer
import numpy as np
import pickle
import os
from docx import Document
from pptx import Presentation

# Function to extract text from docx files
def extract_text_docx(file_path):
    doc = Document(file_path)
    full_text = []
    for para in doc.paragraphs:
        full_text.append(para.text)
    return "\n".join(full_text)

# Function to extract text from pptx files
def extract_text_pptx(file_path):
    prs = Presentation(file_path)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                full_text.append(shape.text)
    return "\n".join(full_text)

# Chunking function remains the same
def chunk_text(text, max_tokens=200):
    sentences = text.split('. ')
    chunks = []
    chunk = ""
    for sentence in sentences:
        if len((chunk + sentence).split()) > max_tokens:
            chunks.append(chunk.strip())
            chunk = sentence + ". "
        else:
            chunk += sentence + ". "
    if chunk:
        chunks.append(chunk.strip())
    return chunks

print("🔄 Fetching documents from Google Drive...")
docs, sources_raw, file_ids_raw, file_paths = fetch_documents(authenticate_drive)

print("✂️ Extracting text, chunking and labeling...")

text_chunks = []
sources = []
file_ids = []

for doc, source, file_id, file_path in zip(docs, sources_raw, file_ids_raw, file_paths):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".docx":
        text = extract_text_docx(file_path)
    elif ext == ".pptx":
        text = extract_text_pptx(file_path)
    else:
        text = doc

    chunks = chunk_text(text)

    # Add file name and folder name as extra searchable chunks
    if " / " in source:
        folder_name, file_name = source.split(" / ", 1)
    else:
        folder_name, file_name = "", source

    chunks.append(file_name)
    if folder_name:
        chunks.append(folder_name)

    text_chunks.extend(chunks)
    sources.extend([source] * len(chunks))
    file_ids.extend([file_id] * len(chunks))

print("🔍 Embedding chunks...")
model = SentenceTransformer('all-MiniLM-L6-v2')
embeddings = model.encode(text_chunks, show_progress_bar=True).astype('float32')

print("💾 Saving data...")
os.makedirs('data', exist_ok=True)
np.save('data/embeddings.npy', embeddings)

with open('data/text_chunks.pkl', 'wb') as f:
    pickle.dump(text_chunks, f)

with open('data/sources.pkl', 'wb') as f:
    pickle.dump(sources, f)

with open('data/file_ids.pkl', 'wb') as f:
    pickle.dump(file_ids, f)

print(" Embedding update complete!")
